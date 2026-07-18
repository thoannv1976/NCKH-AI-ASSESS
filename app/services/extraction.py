"""Trích xuất nội dung văn bản từ sản phẩm nộp (docx/pdf/pptx/xlsx) và liên kết
để đưa vào prompt chấm điểm."""
from __future__ import annotations

import io
from pathlib import Path

# Giới hạn ký tự đưa vào prompt cho mỗi tệp / mỗi phần (an toàn cửa sổ ngữ cảnh)
MAX_CHARS_PER_FILE = 40_000
MAX_CHARS_PER_PART = 120_000


def _truncate(text: str, limit: int) -> str:
    if len(text) <= limit:
        return text
    return text[:limit] + f"\n[... đã cắt bớt, còn {len(text) - limit} ký tự ...]"


def extract_docx(data: bytes) -> str:
    import docx

    doc = docx.Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    pages = []
    for i, page in enumerate(reader.pages[:100]):
        try:
            pages.append(page.extract_text() or "")
        except Exception:  # noqa: BLE001 — trang lỗi thì bỏ qua
            pages.append(f"[Trang {i + 1}: không trích xuất được]")
    text = "\n".join(pages).strip()
    return text or "[PDF không trích xuất được văn bản — có thể là bản scan/ảnh; Hội đồng cần xem bản gốc]"


def extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    out = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
        out.append(f"--- Slide {i} ---\n" + "\n".join(texts))
    out.append(f"[Tổng số slide: {len(prs.slides)}]")
    return "\n".join(out)


def extract_xlsx(data: bytes) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    out = []
    for ws in wb.worksheets:
        out.append(f"--- Sheet: {ws.title} ---")
        for r, row in enumerate(ws.iter_rows(values_only=True)):
            if r > 500:
                out.append("[... cắt bớt các dòng sau ...]")
                break
            cells = ["" if c is None else str(c) for c in row]
            if any(cells):
                out.append(" | ".join(cells))
    return "\n".join(out)


EXTRACTORS = {".docx": extract_docx, ".pdf": extract_pdf, ".pptx": extract_pptx, ".xlsx": extract_xlsx}


def extract_item_text(storage, item: dict) -> str:
    """Trích xuất một submission_item (tệp hoặc liên kết) thành văn bản mô tả + nội dung."""
    name = item.get("original_name") or item.get("url") or "(không tên)"
    if item["type"] == "link":
        status = "truy cập được" if item.get("link_ok") else "KHÔNG truy cập được khi kiểm tra"
        snap = item.get("link_snapshot", "")
        body = f"[Liên kết: {item['url']} — {status}]"
        if snap:
            body += "\nTrích nội dung trang đích:\n" + _truncate(snap, 5_000)
        return body
    ext = Path(name).suffix.lower()
    extractor = EXTRACTORS.get(ext)
    if not extractor:
        return f"[Tệp {name}: định dạng {ext} không trích xuất nội dung được (vd video) — Hội đồng xem trực tiếp]"
    try:
        with storage.open(item["storage_path"]) as f:
            data = f.read()
        return _truncate(extractor(data), MAX_CHARS_PER_FILE)
    except Exception as exc:  # noqa: BLE001 — lỗi trích xuất ghi nhận để chấm thủ công
        return f"[Tệp {name}: lỗi trích xuất nội dung ({type(exc).__name__}) — Hội đồng cần xem bản gốc]"


def build_part_content(storage, items: list[dict], part: str) -> tuple[str, str]:
    """Ghép nội dung (sản phẩm, minh chứng) của một phần thành 2 khối văn bản."""
    def block(kind: str) -> str:
        chunks = []
        for it in [i for i in items if i["part"] == part and i["kind"] == kind]:
            name = it.get("original_name") or it.get("url")
            chunks.append(f"===== {('SẢN PHẨM' if kind == 'product' else 'MINH CHỨNG')}: {name} =====\n"
                          + extract_item_text(storage, it))
        return _truncate("\n\n".join(chunks), MAX_CHARS_PER_PART)

    return block("product"), block("evidence")
